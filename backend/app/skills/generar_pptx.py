"""
Skill: generar_pptx (v2) — Deck comercial profesional SGS España
14 slides con gráficas nativas python-pptx (ColumnChart, BarChart, PieChart),
pipeline activo, análisis de fit, escenarios, ROI, timeline y argumentario.
"""

import io
import json
from datetime import date
from uuid import UUID

import asyncpg
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches, Pt, Emu

# ── Paleta SGS ────────────────────────────────────────────────────────────────
ROJO_SGS    = RGBColor(0xC0, 0x00, 0x1A)
ROJO_CLARO  = RGBColor(0xE8, 0x20, 0x30)
ROJO_MUY_CLARO = RGBColor(0xFF, 0xE5, 0xE8)
NEGRO       = RGBColor(0x0D, 0x0D, 0x0D)
NEGRO_SUAVE = RGBColor(0x1A, 0x1A, 0x2E)
BLANCO      = RGBColor(0xFF, 0xFF, 0xFF)
GRIS        = RGBColor(0x80, 0x80, 0x80)
GRIS_CLARO  = RGBColor(0xF4, 0xF4, 0xF6)
GRIS_MEDIO  = RGBColor(0xCC, 0xCC, 0xD4)
GRIS_OSCURO = RGBColor(0x3A, 0x3A, 0x3A)
VERDE       = RGBColor(0x12, 0x6A, 0x12)
VERDE_CLARO = RGBColor(0x2D, 0xBE, 0x61)
VERDE_BG    = RGBColor(0xE8, 0xF5, 0xE8)
NARANJA     = RGBColor(0xD0, 0x6B, 0x00)
NARANJA_CLARO = RGBColor(0xFF, 0xA5, 0x00)
NARANJA_BG  = RGBColor(0xFF, 0xF3, 0xE0)
AZUL        = RGBColor(0x14, 0x4F, 0xC2)
AZUL_CLARO  = RGBColor(0x3B, 0x82, 0xF6)
AZUL_BG     = RGBColor(0xE8, 0xF0, 0xFA)
AZUL_OSCURO = RGBColor(0x0D, 0x3B, 0x6E)

_PALETTE = [
    RGBColor(0xC0, 0x00, 0x1A),
    RGBColor(0x14, 0x4F, 0xC2),
    RGBColor(0x12, 0x6A, 0x12),
    RGBColor(0xD0, 0x6B, 0x00),
    RGBColor(0x7B, 0x2D, 0xBF),
    RGBColor(0x0A, 0x7A, 0x7A),
    RGBColor(0xC8, 0x5C, 0x10),
    RGBColor(0x2D, 0x6A, 0x9F),
]

ANCHO = Inches(13.33)
ALTO  = Inches(7.5)

_ETAPA_LABEL = {
    "prospection":    "Prospección",
    "qualification":  "Cualificación",
    "proposal":       "Propuesta",
    "negotiation":    "Negociación",
    "value_prop":     "Prop. Valor",
    "id_decision":    "Decisor ID",
    "perception":     "Percepción",
    "closed_won":     "Ganadas",
    "closed_lost":    "Perdidas",
    "closed_withdrawn": "Retiradas",
}


# ── Parsers ───────────────────────────────────────────────────────────────────

def _parse_jsonb(valor, default):
    if isinstance(valor, (dict, list)):
        return valor
    if isinstance(valor, str):
        try:
            return json.loads(valor)
        except (json.JSONDecodeError, ValueError):
            pass
    return default


# ── Presentación base ─────────────────────────────────────────────────────────

def _nueva_presentacion() -> Presentation:
    prs = Presentation()
    prs.slide_width  = ANCHO
    prs.slide_height = ALTO
    return prs


def _slide_en_blanco(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ── Primitivas de dibujo ──────────────────────────────────────────────────────

def _rect(slide, left, top, width, height, color: RGBColor, radio: int = 0):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _texto(slide, texto: str, left, top, width, height,
           size=18, bold=False, color=NEGRO, align=PP_ALIGN.LEFT,
           wrap=True, italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = 1.05
    run = p.add_run()
    run.text = texto
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def _linea(slide, left, top, width, height, color: RGBColor):
    _rect(slide, left, top, width, height, color)


def _bullet_multiline(slide, items: list[str], left, top, width, height,
                      size=12, color=NEGRO, symbol="▸"):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p   = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"{symbol}  {item}"
        run.font.size  = Pt(size)
        run.font.color.rgb = color
    return txb


def _kpi_bloque(slide, left, top, width, height, valor: str, etiqueta: str,
                color_fondo: RGBColor, color_valor: RGBColor):
    _rect(slide, left, top, width, height, color_fondo)
    _texto(slide, valor,
           left, top + Inches(0.15), width, Inches(0.6),
           size=22, bold=True, color=color_valor, align=PP_ALIGN.CENTER)
    _texto(slide, etiqueta,
           left, top + Inches(0.72), width, Inches(0.4),
           size=8, color=GRIS, align=PP_ALIGN.CENTER)


def _barra_progreso(slide, left, top, width, height, pct: float,
                    color_fondo: RGBColor, color_barra: RGBColor):
    _rect(slide, left, top, width, height, color_fondo)
    relleno = int(width * max(0.0, min(1.0, pct / 100)))
    if relleno > 0:
        _rect(slide, left, top, relleno, height, color_barra)


def _color_score(score: float) -> RGBColor:
    if score >= 80:
        return VERDE_CLARO
    if score >= 60:
        return NARANJA_CLARO
    return ROJO_CLARO


def _color_score_bg(score: float) -> RGBColor:
    if score >= 80:
        return VERDE_BG
    if score >= 60:
        return NARANJA_BG
    return ROJO_MUY_CLARO


# ── Header y footer reutilizables ─────────────────────────────────────────────

def _header(slide, titulo: str, subtitulo: str, color_bg=NEGRO_SUAVE):
    titulo_size = 22 if len(titulo) <= 38 else (20 if len(titulo) <= 54 else 18)
    _rect(slide, 0, 0, ANCHO, Inches(1.1), color_bg)
    _rect(slide, 0, 0, Inches(0.06), ALTO, ROJO_SGS)
    _texto(slide, titulo, Inches(0.35), Inches(0.12), Inches(11), Inches(0.7),
           size=titulo_size, bold=True, color=BLANCO, wrap=False)
    if subtitulo:
        _texto(slide, subtitulo, Inches(0.35), Inches(0.77), Inches(11), Inches(0.35),
               size=10, color=RGBColor(0xCC, 0xCC, 0xCC), italic=True)


def _footer(slide, empresa: str = ""):
    _linea(slide, Inches(0.2), Inches(7.2), ANCHO - Inches(0.4), Pt(0.75), GRIS_MEDIO)
    fecha = date.today().strftime("%B %Y").capitalize()
    _texto(slide, f"SGS España · Inteligencia Comercial · {fecha} · Confidencial",
           Inches(0.3), Inches(7.25), Inches(8), Inches(0.3), size=8, color=GRIS)
    if empresa:
        _texto(slide, empresa, Inches(9), Inches(7.25), Inches(4), Inches(0.3),
               size=8, color=GRIS, align=PP_ALIGN.RIGHT)


# ── Gráficas nativas python-pptx ──────────────────────────────────────────────

def _chart_column(slide, left, top, width, height,
                  categories: list, values: list,
                  serie_label: str = "", colores: list | None = None):
    """BarChart vertical nativo."""
    cd = ChartData()
    cd.categories = categories
    cd.add_series(serie_label, [float(v) for v in values])

    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, cd
    )
    chart = gf.chart
    chart.has_legend = False
    chart.has_title  = False

    # Color de las barras
    series = chart.series[0]
    try:
        for i, point in enumerate(series.points):
            c = colores[i % len(colores)] if colores else _PALETTE[i % len(_PALETTE)]
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = c
    except Exception:
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = ROJO_SGS

    # Etiquetas de valor
    try:
        plot = chart.plots[0]
        plot.has_data_labels = True
        dls = plot.data_labels
        dls.font.size = Pt(8)
        dls.font.bold = True
    except Exception:
        pass

    # Quitar borde exterior del chart
    try:
        gf.line.color.type = None
    except Exception:
        pass

    return gf


def _chart_bar_horizontal(slide, left, top, width, height,
                           categories: list, values: list,
                           colores: list | None = None):
    """BarChart horizontal nativo."""
    cd = ChartData()
    cd.categories = categories
    cd.add_series("", [float(v) for v in values])

    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, cd
    )
    chart = gf.chart
    chart.has_legend = False
    chart.has_title  = False

    series = chart.series[0]
    try:
        for i, point in enumerate(series.points):
            c = colores[i % len(colores)] if colores else _PALETTE[i % len(_PALETTE)]
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = c
    except Exception:
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = ROJO_SGS

    try:
        chart.plots[0].has_data_labels = True
        chart.plots[0].data_labels.font.size = Pt(8)
    except Exception:
        pass

    return gf


def _chart_pie(slide, left, top, width, height,
               categories: list, values: list):
    """PieChart nativo con colores SGS."""
    cd = ChartData()
    cd.categories = categories
    cd.add_series("", [float(v) for v in values])

    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, left, top, width, height, cd
    )
    chart = gf.chart
    chart.has_legend = True
    chart.has_title  = False

    try:
        for i, point in enumerate(chart.series[0].points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = _PALETTE[i % len(_PALETTE)]
    except Exception:
        pass

    try:
        plot = chart.plots[0]
        plot.has_data_labels = True
        dls = plot.data_labels
        dls.number_format = "0%"
        dls.font.size = Pt(8)
        dls.font.bold = True
        dls.font.color.rgb = BLANCO
    except Exception:
        pass

    return gf


# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════

# ── Slide 1: Portada ──────────────────────────────────────────────────────────

def _slide_portada(prs, nombre: str, sector: str | None,
                   opt: dict, med: dict, pes: dict,
                   num_productos: int, total_ops: int, total_pipeline: float):
    slide = _slide_en_blanco(prs)

    # Fondo oscuro completo
    _rect(slide, 0, 0, ANCHO, ALTO, NEGRO_SUAVE)

    # Banda roja izquierda gruesa
    _rect(slide, 0, 0, Inches(0.35), ALTO, ROJO_SGS)

    # Banda roja superior fina
    _rect(slide, Inches(0.35), 0, ANCHO - Inches(0.35), Inches(0.06), ROJO_SGS)

    # Marca de agua decorativa (rectángulo grande translúcido simulado con color próximo)
    _rect(slide, Inches(9), Inches(1.5), Inches(5), Inches(5), RGBColor(0x22, 0x22, 0x40))

    # Logo / Brand
    _texto(slide, "SGS España", Inches(0.7), Inches(0.4), Inches(6), Inches(0.85),
           size=36, bold=True, color=BLANCO)
    _texto(slide, "Sistema de Inteligencia Comercial",
           Inches(0.7), Inches(1.2), Inches(7), Inches(0.45),
           size=13, color=GRIS, italic=True)

    # Línea divisora roja
    _linea(slide, Inches(0.7), Inches(1.75), Inches(5), Pt(2), ROJO_SGS)

    # Tipo de documento
    _texto(slide, "PROPUESTA COMERCIAL ESTRATÉGICA",
           Inches(0.7), Inches(2.0), Inches(9), Inches(0.5),
           size=11, color=RGBColor(0xAA, 0xAA, 0xBB), bold=True)

    # Nombre empresa — protagonista
    _texto(slide, nombre, Inches(0.7), Inches(2.5), Inches(9.5), Inches(1.5),
           size=40, bold=True, color=BLANCO)

    # Sector
    if sector:
        _texto(slide, sector.upper(), Inches(0.7), Inches(4.05), Inches(7), Inches(0.45),
               size=12, color=ROJO_CLARO, bold=True)

    # ── KPIs en fila (4 bloques) ──────────────────────────────────────────────
    kpi_top  = Inches(4.75)
    kpi_h    = Inches(1.5)
    kpi_w    = Inches(2.9)
    gap      = Inches(0.18)
    kpi_left = Inches(0.7)

    importe_opt = float(opt.get("importe", 0) or 0)
    importe_med = float(med.get("importe", 0) or 0)
    prob_med    = med.get("probabilidad", 0)

    _kpi_bloque(slide, kpi_left, kpi_top, kpi_w, kpi_h,
                f"{importe_opt:,.0f} €", "Escenario Optimista", ROJO_SGS, BLANCO)
    _kpi_bloque(slide, kpi_left + kpi_w + gap, kpi_top, kpi_w, kpi_h,
                f"{importe_med:,.0f} €", "Escenario Recomendado",
                RGBColor(0x25, 0x25, 0x45), BLANCO)
    _kpi_bloque(slide, kpi_left + (kpi_w + gap) * 2, kpi_top, kpi_w, kpi_h,
                f"{total_pipeline:,.0f} €", f"Pipeline Activo ({total_ops} ops)",
                RGBColor(0x25, 0x25, 0x45), RGBColor(0x6B, 0xC8, 0xFF))
    _kpi_bloque(slide, kpi_left + (kpi_w + gap) * 3, kpi_top, kpi_w, kpi_h,
                f"{prob_med}%", "Probabilidad Media",
                RGBColor(0x25, 0x25, 0x45), RGBColor(0x6B, 0xFF, 0xAA))

    # Fecha y confidencial
    _texto(slide, date.today().strftime("%B de %Y").capitalize(),
           Inches(0.7), Inches(7.1), Inches(4), Inches(0.35),
           size=9, color=GRIS)
    _texto(slide, "CONFIDENCIAL · Uso exclusivo del receptor",
           Inches(7), Inches(7.1), Inches(6), Inches(0.35),
           size=9, color=GRIS, align=PP_ALIGN.RIGHT)


# ── Slide 2: Agenda ───────────────────────────────────────────────────────────

def _slide_agenda(prs, nombre: str):
    slide = _slide_en_blanco(prs)
    _rect(slide, 0, 0, ANCHO, ALTO, BLANCO)
    _header(slide, "Agenda", f"Contenido de la propuesta para {nombre}", ROJO_SGS)

    items = [
        ("01", "Contexto empresarial", "Ficha, certificaciones y presencia digital"),
        ("02", "Diagnóstico de necesidades", "Pain points identificados y oportunidades de mejora"),
        ("03", "Pipeline y actividad comercial", "Análisis de oportunidades activas por etapa e importe"),
        ("04", "Análisis de fit — Productos SGS", "Evaluación de adecuación por producto con score de afinidad"),
        ("05", "Escenarios de venta", "Comparativa pesimista / recomendado / optimista con importes y plazos"),
        ("06", "ROI y proyección de valor", "Retorno esperado y análisis de valor comercial"),
        ("07", "Plan de acción y timeline", "Acciones priorizadas con plazos y responsables"),
        ("08", "Próximos pasos", "Compromisos y hoja de ruta acordada"),
    ]

    cols = [items[:4], items[4:]]
    for col_idx, col_items in enumerate(cols):
        x_base = Inches(0.5 + col_idx * 6.5)
        for i, (num, titulo, desc) in enumerate(col_items):
            y = Inches(1.3 + i * 1.45)
            _rect(slide, x_base, y, Inches(0.55), Inches(0.55), ROJO_SGS)
            _texto(slide, num, x_base, y - Pt(2), Inches(0.55), Inches(0.6),
                   size=13, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
            _texto(slide, titulo, x_base + Inches(0.7), y, Inches(5.6), Inches(0.35),
                   size=13, bold=True, color=NEGRO)
            _texto(slide, desc, x_base + Inches(0.7), y + Inches(0.38), Inches(5.6), Inches(0.5),
                   size=9, color=GRIS, italic=True)
            _linea(slide, x_base, y + Inches(1.3), Inches(6.2), Pt(0.5), GRIS_MEDIO)

    _footer(slide, nombre)


# ── Slide 3: Ficha empresa ────────────────────────────────────────────────────

def _slide_ficha(prs, datos: dict, nombre: str):
    slide = _slide_en_blanco(prs)
    _rect(slide, 0, 0, ANCHO, ALTO, BLANCO)
    _header(slide, "Contexto Empresarial",
            "Ficha de empresa, presencia digital y certificaciones actuales")

    campos = [
        ("Razón social", nombre),
        ("Sector de actividad", datos.get("sector") or "No identificado"),
        ("Número de empleados", str(datos.get("num_empleados") or "No disponible")),
        ("Facturación estimada", datos.get("facturacion_estimada") or "No disponible"),
        ("Presencia web", (datos.get("presencia_web") or "No encontrada")[:60]),
    ]

    # Columna izquierda: datos empresa
    _rect(slide, Inches(0.2), Inches(1.2), Inches(6.2), ALTO - Inches(1.4), GRIS_CLARO)
    _texto(slide, "DATOS DE LA EMPRESA", Inches(0.35), Inches(1.28), Inches(5.5), Inches(0.4),
           size=9, bold=True, color=ROJO_SGS)
    _linea(slide, Inches(0.35), Inches(1.65), Inches(5.8), Pt(1), ROJO_SGS)

    for i, (campo, valor) in enumerate(campos):
        y   = Inches(1.8 + i * 0.95)
        bg  = BLANCO if i % 2 == 0 else GRIS_CLARO
        _rect(slide, Inches(0.25), y, Inches(6.1), Inches(0.82), bg)
        _texto(slide, campo, Inches(0.4), y + Pt(5), Inches(2.3), Inches(0.55),
               size=9, bold=True, color=GRIS_OSCURO)
        _texto(slide, str(valor), Inches(2.8), y + Pt(5), Inches(3.5), Inches(0.55),
               size=10, color=NEGRO)

    # Columna derecha: certificaciones + presencia
    certs = _parse_jsonb(datos.get("certificaciones_actuales"), [])
    _rect(slide, Inches(6.6), Inches(1.2), Inches(6.4), ALTO - Inches(1.4), GRIS_CLARO)
    _texto(slide, "CERTIFICACIONES ACTUALES", Inches(6.75), Inches(1.28), Inches(6), Inches(0.4),
           size=9, bold=True, color=ROJO_SGS)
    _linea(slide, Inches(6.75), Inches(1.65), Inches(6.0), Pt(1), ROJO_SGS)

    if certs:
        for i, cert in enumerate(certs[:6]):
            y = Inches(1.8 + i * 0.7)
            _rect(slide, Inches(6.75), y + Pt(5), Inches(0.22), Inches(0.22), VERDE_CLARO)
            _texto(slide, str(cert), Inches(7.1), y, Inches(5.7), Inches(0.6),
                   size=11, color=NEGRO)
    else:
        _rect(slide, Inches(6.75), Inches(1.85), Inches(5.9), Inches(0.7), NARANJA_BG)
        _texto(slide, "Sin certificaciones detectadas — oportunidad de entrada SGS",
               Inches(6.9), Inches(1.95), Inches(5.6), Inches(0.5),
               size=10, color=NARANJA, italic=True)

    # Investigación IA
    inv_fecha = str(datos.get("investigacion_fecha") or "")[:10]
    modelo    = datos.get("modelo_usado") or "Sistema IC SGS"
    _rect(slide, Inches(6.75), Inches(6.0), Inches(5.9), Inches(0.9), AZUL_BG)
    _texto(slide, "Investigación IA:",
           Inches(6.9), Inches(6.05), Inches(2.5), Inches(0.4), size=9, bold=True, color=AZUL)
    _texto(slide, f"Fecha: {inv_fecha}  ·  Modelo: {modelo}",
           Inches(6.9), Inches(6.42), Inches(5.6), Inches(0.35), size=9, color=AZUL_OSCURO)

    _footer(slide, nombre)


# ── Slide 4: Pain Points + Oportunidades IA ──────────────────────────────────

def _slide_diagnostico(prs, pain: list, oportunidades_ia: list, nombre: str):
    slide = _slide_en_blanco(prs)
    _rect(slide, 0, 0, ANCHO, ALTO, BLANCO)
    _header(slide, "Diagnóstico de Necesidades",
            "Necesidades críticas identificadas y oportunidades de mejora detectadas por IA")

    # Panel izquierdo — Pain Points
    _rect(slide, Inches(0.2), Inches(1.2), Inches(6.2), ALTO - Inches(1.4), ROJO_MUY_CLARO)
    _texto(slide, "PAIN POINTS IDENTIFICADOS", Inches(0.35), Inches(1.28), Inches(5.8), Inches(0.4),
           size=9, bold=True, color=ROJO_SGS)
    _linea(slide, Inches(0.35), Inches(1.65), Inches(5.8), Pt(1.5), ROJO_SGS)

    for i, pp in enumerate(pain[:6]):
        y   = Inches(1.8 + i * 0.88)
        prioridad = "ALTA" if i < 2 else ("MEDIA" if i < 4 else "BAJA")
        c_pill    = ROJO_SGS if i < 2 else (NARANJA if i < 4 else VERDE)
        _rect(slide, Inches(0.3), y + Pt(4), Inches(0.4), Pt(18), c_pill)
        _texto(slide, str(i + 1), Inches(0.3), y + Pt(2), Inches(0.4), Pt(22),
               size=10, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
        _texto(slide, str(pp)[:110], Inches(0.82), y, Inches(5.3), Inches(0.75),
               size=10, color=NEGRO)
        _texto(slide, prioridad, Inches(5.3), y + Pt(4), Inches(0.9), Pt(18),
               size=7, bold=True, color=c_pill, align=PP_ALIGN.RIGHT)
        _linea(slide, Inches(0.35), y + Inches(0.78), Inches(5.8), Pt(0.5), GRIS_MEDIO)

    if not pain:
        _texto(slide, "No se han identificado pain points específicos.",
               Inches(0.5), Inches(2.2), Inches(5.8), Inches(0.5), size=10, color=GRIS, italic=True)

    # Panel derecho — Oportunidades IA
    _rect(slide, Inches(6.6), Inches(1.2), Inches(6.4), ALTO - Inches(1.4), VERDE_BG)
    _texto(slide, "OPORTUNIDADES DE NEGOCIO DETECTADAS", Inches(6.75), Inches(1.28), Inches(6), Inches(0.4),
           size=9, bold=True, color=VERDE)
    _linea(slide, Inches(6.75), Inches(1.65), Inches(6.0), Pt(1.5), VERDE)

    for i, op in enumerate((oportunidades_ia or [])[:6]):
        y = Inches(1.8 + i * 0.88)
        _rect(slide, Inches(6.75), y + Pt(4), Pt(18), Pt(18), VERDE_CLARO)
        _texto(slide, "✓", Inches(6.75), y + Pt(2), Pt(22), Pt(22),
               size=9, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
        txt = op if isinstance(op, str) else str(op)
        _texto(slide, txt[:110], Inches(7.25), y, Inches(5.5), Inches(0.75),
               size=10, color=NEGRO)
        _linea(slide, Inches(6.75), y + Inches(0.78), Inches(6.0), Pt(0.5), GRIS_MEDIO)

    if not oportunidades_ia:
        _texto(slide, "No se han identificado oportunidades específicas.",
               Inches(6.9), Inches(2.2), Inches(5.8), Inches(0.5), size=10, color=GRIS, italic=True)

    _footer(slide, nombre)


# ── Slide 5: Pipeline activo (ColumnChart nativo) ─────────────────────────────

def _slide_pipeline(prs, ops_rows: list, total_ops: int,
                    total_pipeline: float, nombre: str):
    slide = _slide_en_blanco(prs)
    _rect(slide, 0, 0, ANCHO, ALTO, BLANCO)
    _header(slide, "Pipeline y Actividad Comercial",
            "Distribución del pipeline activo por etapa — oportunidades y volumen económico")

    # KPIs rápidos
    kpi_w = Inches(2.8)
    kpi_h = Inches(0.85)
    kpi_y = Inches(1.25)
    gap   = Inches(0.18)

    avg = total_pipeline / total_ops if total_ops else 0
    _kpi_bloque(slide, Inches(0.2), kpi_y, kpi_w, kpi_h,
                f"{total_ops}", "Oportunidades activas", AZUL_BG, AZUL)
    _kpi_bloque(slide, Inches(0.2) + kpi_w + gap, kpi_y, kpi_w, kpi_h,
                f"{total_pipeline:,.0f} €", "Pipeline total", VERDE_BG, VERDE)
    _kpi_bloque(slide, Inches(0.2) + (kpi_w + gap) * 2, kpi_y, kpi_w, kpi_h,
                f"{len(ops_rows)}", "Etapas con actividad", GRIS_CLARO, GRIS_OSCURO)
    _kpi_bloque(slide, Inches(0.2) + (kpi_w + gap) * 3, kpi_y, kpi_w, kpi_h,
                f"{avg:,.0f} €", "Importe medio/op", NARANJA_BG, NARANJA)

    if ops_rows:
        etapas  = [_ETAPA_LABEL.get(r["etapa"], r["etapa"].replace("_", " ").title())
                   for r in ops_rows[:8]]
        importes = [float(r["importe_total"]) for r in ops_rows[:8]]

        # Gráfica nativa — columna izquierda
        _chart_column(
            slide,
            Inches(0.2), Inches(2.3), Inches(7.8), Inches(4.5),
            etapas, importes,
        )

        # Tabla resumen — columna derecha
        _rect(slide, Inches(8.3), Inches(2.2), Inches(4.8), ALTO - Inches(2.4), GRIS_CLARO)
        _texto(slide, "DETALLE POR ETAPA", Inches(8.45), Inches(2.3), Inches(4.5), Inches(0.35),
               size=9, bold=True, color=GRIS_OSCURO)
        _linea(slide, Inches(8.45), Inches(2.62), Inches(4.6), Pt(0.75), GRIS_MEDIO)

        # Cabecera tabla
        _rect(slide, Inches(8.3), Inches(2.68), Inches(4.8), Inches(0.4), NEGRO_SUAVE)
        _texto(slide, "Etapa", Inches(8.42), Inches(2.72), Inches(2.4), Inches(0.35),
               size=8, bold=True, color=BLANCO)
        _texto(slide, "Ops", Inches(10.9), Inches(2.72), Inches(0.7), Inches(0.35),
               size=8, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
        _texto(slide, "Importe", Inches(11.65), Inches(2.72), Inches(1.3), Inches(0.35),
               size=8, bold=True, color=BLANCO, align=PP_ALIGN.RIGHT)

        for i, r in enumerate(ops_rows[:7]):
            y   = Inches(3.15 + i * 0.52)
            bg  = BLANCO if i % 2 == 0 else GRIS_CLARO
            _rect(slide, Inches(8.3), y, Inches(4.8), Inches(0.48), bg)
            lbl = _ETAPA_LABEL.get(r["etapa"], r["etapa"].replace("_", " ").title())
            _texto(slide, lbl, Inches(8.42), y + Pt(4), Inches(2.4), Inches(0.38),
                   size=9, color=NEGRO)
            _texto(slide, str(r["total_ops"]), Inches(10.9), y + Pt(4), Inches(0.7), Inches(0.38),
                   size=9, bold=True, color=AZUL, align=PP_ALIGN.CENTER)
            imp = float(r["importe_total"])
            _texto(slide, f"{imp:,.0f} €", Inches(11.65), y + Pt(4), Inches(1.3), Inches(0.38),
                   size=9, bold=True, color=NEGRO, align=PP_ALIGN.RIGHT)

    else:
        _rect(slide, Inches(0.5), Inches(2.5), Inches(12), Inches(1.2), NARANJA_BG)
        _texto(slide, "No hay oportunidades activas en pipeline — oportunidad de apertura para productos SGS.",
               Inches(0.7), Inches(2.75), Inches(11), Inches(0.7), size=13, color=NARANJA, italic=True)

    _footer(slide, nombre)


# ── Slide 6: Distribución por producto (PieChart) ─────────────────────────────

def _slide_distribucion_productos(prs, prods_pipeline: list, total_pipeline: float, nombre: str):
    slide = _slide_en_blanco(prs)
    _rect(slide, 0, 0, ANCHO, ALTO, BLANCO)
    _header(slide, "Distribución del Pipeline por Producto",
            "Concentración del volumen de oportunidades activas por línea de negocio")

    if prods_pipeline:
        nombres_pie  = [str(r["producto"])[:25] for r in prods_pipeline[:6]]
        valores_pie  = [float(r["importe"]) for r in prods_pipeline[:6]]
        total_mostrado = sum(valores_pie)

        # PieChart nativo izquierda
        _chart_pie(
            slide,
            Inches(0.2), Inches(1.2), Inches(6.8), Inches(5.5),
            nombres_pie, valores_pie,
        )

        # Tabla derecha
        _rect(slide, Inches(7.3), Inches(1.2), Inches(5.8), ALTO - Inches(1.4), GRIS_CLARO)
        _texto(slide, "VOLUMEN POR PRODUCTO", Inches(7.45), Inches(1.3), Inches(5.5), Inches(0.4),
               size=9, bold=True, color=GRIS_OSCURO)
        _linea(slide, Inches(7.45), Inches(1.68), Inches(5.6), Pt(0.75), GRIS_MEDIO)

        # Cabecera
        _rect(slide, Inches(7.3), Inches(1.75), Inches(5.8), Inches(0.42), NEGRO_SUAVE)
        _texto(slide, "Producto", Inches(7.45), Inches(1.79), Inches(2.8), Inches(0.35),
               size=8, bold=True, color=BLANCO)
        _texto(slide, "Ops", Inches(10.3), Inches(1.79), Inches(0.7), Inches(0.35),
               size=8, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
        _texto(slide, "Importe", Inches(11.1), Inches(1.79), Inches(1.8), Inches(0.35),
               size=8, bold=True, color=BLANCO, align=PP_ALIGN.RIGHT)

        for i, r in enumerate(prods_pipeline[:7]):
            y  = Inches(2.23 + i * 0.63)
            bg = BLANCO if i % 2 == 0 else GRIS_CLARO
            _rect(slide, Inches(7.3), y, Inches(5.8), Inches(0.57), bg)
            _rect(slide, Inches(7.35), y + Pt(10), Pt(10), Pt(10),
                  _PALETTE[i % len(_PALETTE)])
            _texto(slide, str(r["producto"])[:30], Inches(7.65), y + Pt(5), Inches(2.65), Inches(0.45),
                   size=10, color=NEGRO)
            _texto(slide, str(r["total"]), Inches(10.3), y + Pt(5), Inches(0.7), Inches(0.45),
                   size=10, bold=True, color=AZUL, align=PP_ALIGN.CENTER)
            pct = (float(r["importe"]) / total_pipeline * 100) if total_pipeline else 0
            _texto(slide, f"{float(r['importe']):,.0f} € ({pct:.0f}%)",
                   Inches(11.1), y + Pt(5), Inches(1.9), Inches(0.45),
                   size=9, bold=True, color=NEGRO, align=PP_ALIGN.RIGHT)

    else:
        _texto(slide, "Sin datos de pipeline por producto.",
               Inches(2), Inches(3), Inches(9), Inches(1),
               size=14, color=GRIS, italic=True, align=PP_ALIGN.CENTER)

    _footer(slide, nombre)


# ── Slide 7: Análisis de fit productos SGS ────────────────────────────────────

def _slide_fit_productos(prs, productos: list, nombre: str):
    slide = _slide_en_blanco(prs)
    _rect(slide, 0, 0, ANCHO, ALTO, BLANCO)
    _header(slide, "Análisis de Fit — Productos SGS Recomendados",
            "Score de adecuación ponderado por sector, pain points, certificaciones y pipeline histórico")

    # Cabecera tabla
    _rect(slide, Inches(0.2), Inches(1.2), ANCHO - Inches(0.4), Inches(0.42), NEGRO_SUAVE)
    _texto(slide, "Producto / Servicio SGS", Inches(0.35), Inches(1.25), Inches(3.8), Inches(0.35),
           size=9, bold=True, color=BLANCO)
    _texto(slide, "Score", Inches(4.25), Inches(1.25), Inches(1.0), Inches(0.35),
           size=9, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    _texto(slide, "Indicador visual", Inches(5.4), Inches(1.25), Inches(2.8), Inches(0.35),
           size=9, bold=True, color=BLANCO)
    _texto(slide, "Argumentario comercial clave", Inches(8.35), Inches(1.25), Inches(4.8), Inches(0.35),
           size=9, bold=True, color=BLANCO)

    prods_sorted = sorted(productos, key=lambda x: -x.get("score_fit", 0))[:7]
    for i, prod in enumerate(prods_sorted):
        y     = Inches(1.7 + i * 0.8)
        bg    = GRIS_CLARO if i % 2 == 0 else BLANCO
        score = float(prod.get("score_fit", 0))
        c_s   = _color_score(score)
        c_bg  = _color_score_bg(score)

        _rect(slide, Inches(0.2), y, ANCHO - Inches(0.4), Inches(0.72), bg)

        # Producto
        _texto(slide, str(prod.get("producto", ""))[:38],
               Inches(0.35), y + Pt(6), Inches(3.8), Inches(0.55),
               size=11, bold=True, color=NEGRO)

        # Score numérico con fondo de color
        _rect(slide, Inches(4.25), y + Pt(6), Inches(1.0), Inches(0.42), c_bg)
        _texto(slide, f"{score:.0f}%", Inches(4.25), y + Pt(5), Inches(1.0), Inches(0.45),
               size=14, bold=True, color=c_s, align=PP_ALIGN.CENTER)

        # Barra progreso
        _barra_progreso(slide, Inches(5.4), y + Pt(14), Inches(2.7), Pt(14),
                        score, GRIS_MEDIO, c_s)

        # Argumentario
        arg = str(prod.get("argumentario", ""))[:130]
        _texto(slide, arg, Inches(8.35), y + Pt(5), Inches(4.8), Inches(0.6),
               size=8.5, color=GRIS_OSCURO)

    # Leyenda
    _texto(slide, "● ≥80% Alta afinidad    ● ≥60% Media    ● <60% Baja",
           Inches(0.3), Inches(7.22), Inches(8), Inches(0.28), size=8, color=GRIS)
    _footer(slide, nombre)


# ── Slide 8: Escenarios de venta (ColumnChart nativo) ────────────────────────

def _slide_escenarios(prs, opt: dict, med: dict, pes: dict, nombre: str):
    slide = _slide_en_blanco(prs)
    _rect(slide, 0, 0, ANCHO, ALTO, BLANCO)
    _header(slide, "Escenarios de Venta — Análisis Económico Comparado",
            "Proyección de cierre por escenario basada en pipeline actual y probabilidad histórica")

    imp_pes = float(pes.get("importe", 0) or 0)
    imp_med = float(med.get("importe", 0) or 0)
    imp_opt = float(opt.get("importe", 0) or 0)
    prob_pes = pes.get("probabilidad", 0)
    prob_med = med.get("probabilidad", 0)
    prob_opt = opt.get("probabilidad", 0)
    pond_pes = imp_pes * prob_pes / 100
    pond_med = imp_med * prob_med / 100
    pond_opt = imp_opt * prob_opt / 100

    # Gráfica comparativa — columna izquierda
    _chart_column(
        slide,
        Inches(0.2), Inches(1.25), Inches(6.8), Inches(4.5),
        ["Pesimista", "Recomendado", "Optimista"],
        [imp_pes, imp_med, imp_opt],
        colores=[ROJO_CLARO, NARANJA_CLARO, VERDE_CLARO],
    )

    # Tabla de detalle — columna derecha
    escenarios = [
        ("Pesimista",   pes, ROJO_MUY_CLARO,  ROJO_SGS),
        ("Recomendado", med, NARANJA_BG,       NARANJA),
        ("Optimista",   opt, VERDE_BG,         VERDE),
    ]
    for i, (nombre_esc, esc, bg, c) in enumerate(escenarios):
        y = Inches(1.25 + i * 1.95)
        _rect(slide, Inches(7.2), y, Inches(5.9), Inches(1.8), bg)
        _linea(slide, Inches(7.2), y, Inches(0.12), Inches(1.8), c)
        _texto(slide, nombre_esc.upper(), Inches(7.45), y + Pt(8), Inches(3), Inches(0.38),
               size=10, bold=True, color=c)
        _texto(slide, f"{float(esc.get('importe', 0)):,.0f} €",
               Inches(7.45), y + Pt(28), Inches(3.5), Inches(0.55),
               size=22, bold=True, color=c)
        _texto(slide, f"Probabilidad: {esc.get('probabilidad', 0)}%  ·  Plazo: {esc.get('plazo_meses', '?')} meses",
               Inches(7.45), y + Pt(68), Inches(5.5), Inches(0.35),
               size=9, color=GRIS_OSCURO)

    # Importe ponderado (esperanza matemática)
    _rect(slide, Inches(0.2), Inches(5.95), ANCHO - Inches(0.4), Inches(0.75), NEGRO_SUAVE)
    _texto(slide, "Importe ponderado (esperanza de ingreso = importe × probabilidad):",
           Inches(0.4), Inches(6.03), Inches(5.5), Inches(0.35), size=9, color=GRIS, italic=True)
    for i, (lbl, pond, c) in enumerate([
        ("Pesimista", pond_pes, ROJO_CLARO),
        ("Recomendado", pond_med, NARANJA_CLARO),
        ("Optimista", pond_opt, VERDE_CLARO),
    ]):
        x = Inches(6.0 + i * 2.3)
        _texto(slide, f"{pond:,.0f} €", x, Inches(6.0), Inches(2.2), Inches(0.4),
               size=14, bold=True, color=c, align=PP_ALIGN.CENTER)
        _texto(slide, lbl, x, Inches(6.38), Inches(2.2), Inches(0.28),
               size=8, color=GRIS, align=PP_ALIGN.CENTER)

    _footer(slide, nombre)


# ── Slide 9: ROI y Proyección de Valor ───────────────────────────────────────

def _slide_roi(prs, opt: dict, med: dict, pes: dict, num_productos: int, nombre: str):
    slide = _slide_en_blanco(prs)
    _rect(slide, 0, 0, ANCHO, ALTO, BLANCO)
    _header(slide, "ROI y Proyección de Valor Comercial",
            "Análisis de retorno esperado y valor estratégico del acuerdo", AZUL_OSCURO)

    imp_med  = float(med.get("importe", 0) or 0)
    imp_opt  = float(opt.get("importe", 0) or 0)
    imp_pes  = float(pes.get("importe", 0) or 0)
    prob_med = med.get("probabilidad", 0)
    prob_opt = opt.get("probabilidad", 0)
    plazo_m  = med.get("plazo_meses", 12)

    # Ingreso ponderado esperado
    esperanza      = imp_med * prob_med / 100
    # Rango de valor
    rango_bajo     = imp_pes
    rango_alto     = imp_opt
    # Valor anualizado (estimación)
    meses          = max(int(plazo_m or 12), 1)
    valor_mes      = imp_med / meses
    valor_anual    = valor_mes * 12

    metricas = [
        ("Valor esperado ponderado",    f"{esperanza:,.0f} €",         AZUL, AZUL_BG,    "Importe medio × probabilidad"),
        ("Rango de valor del acuerdo",  f"{rango_bajo:,.0f} – {rango_alto:,.0f} €", VERDE, VERDE_BG, "Entre escenario pesimista y optimista"),
        ("Importe anualizado est.",      f"{valor_anual:,.0f} €/año",   NARANJA, NARANJA_BG, f"Basado en plazo {meses} meses"),
        ("Nº productos involucrados",   str(num_productos),             ROJO_SGS, ROJO_MUY_CLARO, "Líneas de negocio SGS con alta afinidad"),
        ("Ciclo de venta estimado",     f"{meses} meses",               AZUL_OSCURO, AZUL_BG, "Escenario recomendado"),
    ]

    for i, (metrica, valor, c_val, c_bg, nota) in enumerate(metricas):
        y  = Inches(1.3 + i * 1.08)
        _rect(slide, Inches(0.3), y, Inches(12.6), Inches(0.96), c_bg)
        _linea(slide, Inches(0.3), y, Inches(0.15), Inches(0.96), c_val)
        _texto(slide, metrica, Inches(0.65), y + Pt(6), Inches(5.5), Inches(0.4),
               size=12, bold=True, color=NEGRO)
        _texto(slide, nota, Inches(0.65), y + Pt(40), Inches(5.5), Inches(0.32),
               size=8, color=GRIS, italic=True)
        _texto(slide, valor, Inches(6.5), y + Pt(10), Inches(6.3), Inches(0.5),
               size=20, bold=True, color=c_val, align=PP_ALIGN.RIGHT)

    _footer(slide, nombre)


# ── Slide 10: Análisis de fit con BarChart horizontal ─────────────────────────

def _slide_fit_chart(prs, productos: list, nombre: str):
    slide = _slide_en_blanco(prs)
    _rect(slide, 0, 0, ANCHO, ALTO, BLANCO)
    _header(slide, "Score de Fit — Comparativa Visual",
            "Ranking de productos SGS ordenado por score de adecuación")

    prods_sorted = sorted(productos, key=lambda x: -x.get("score_fit", 0))[:8]
    nombres_chart = [str(p.get("producto", ""))[:30] for p in prods_sorted]
    scores_chart  = [float(p.get("score_fit", 0)) for p in prods_sorted]
    colores_chart = [_color_score(s) for s in scores_chart]

    # BarChart horizontal nativo
    _chart_bar_horizontal(
        slide,
        Inches(0.2), Inches(1.2), Inches(9.5), Inches(5.7),
        nombres_chart, scores_chart,
        colores=colores_chart,
    )

    # Leyenda derecha
    _rect(slide, Inches(10.0), Inches(1.2), Inches(3.1), Inches(5.7), GRIS_CLARO)
    _texto(slide, "LEYENDA DE SCORES", Inches(10.15), Inches(1.3), Inches(2.8), Inches(0.35),
           size=9, bold=True, color=GRIS_OSCURO)
    _linea(slide, Inches(10.15), Inches(1.63), Inches(2.8), Pt(0.75), GRIS_MEDIO)
    for top_pct, lbl, c, bg in [
        (80, "≥ 80% — Alta afinidad", VERDE, VERDE_BG),
        (60, "≥ 60% — Afinidad media", NARANJA, NARANJA_BG),
        (0,  "< 60% — Afinidad baja", ROJO_SGS, ROJO_MUY_CLARO),
    ]:
        _rect(slide, Inches(10.15), Inches(1.82 + (80 - top_pct) / 80 * 2.5), Inches(2.8), Inches(0.72), bg)
        _texto(slide, lbl, Inches(10.3), Inches(1.95 + (80 - top_pct) / 80 * 2.5), Inches(2.6), Inches(0.35),
               size=10, bold=True, color=c)

    _footer(slide, nombre)


# ── Slide 11: Plan de acción (tabla) ──────────────────────────────────────────

def _slide_plan_accion(prs, plan: list, nombre: str):
    slide = _slide_en_blanco(prs)
    _rect(slide, 0, 0, ANCHO, ALTO, BLANCO)
    _header(slide, "Plan de Acción Priorizado",
            "Acciones ordenadas por prioridad comercial con tipo, plazo e impacto esperado")

    tipo_c = {
        "nuevo":      VERDE,
        "renovacion": NARANJA,
        "upselling":  AZUL,
    }
    tipo_bg = {
        "nuevo":      VERDE_BG,
        "renovacion": NARANJA_BG,
        "upselling":  AZUL_BG,
    }

    # Cabecera
    _rect(slide, Inches(0.2), Inches(1.2), ANCHO - Inches(0.4), Inches(0.42), NEGRO_SUAVE)
    for txt, x, w, aln in [
        ("#",       0.22, 0.48, PP_ALIGN.CENTER),
        ("ACCIÓN",  0.82, 7.0,  PP_ALIGN.LEFT),
        ("TIPO",    8.0,  1.6,  PP_ALIGN.CENTER),
        ("PLAZO",   9.75, 1.0,  PP_ALIGN.CENTER),
        ("IMPACTO", 10.9, 2.1,  PP_ALIGN.LEFT),
    ]:
        _texto(slide, txt, Inches(x), Inches(1.25), Inches(w), Inches(0.35),
               size=9, bold=True, color=BLANCO, align=aln)

    acciones = sorted(plan, key=lambda x: x.get("prioridad", 99))[:7]
    for i, acc in enumerate(acciones):
        y     = Inches(1.7 + i * 0.77)
        bg    = GRIS_CLARO if i % 2 == 0 else BLANCO
        tipo  = str(acc.get("tipo", "nuevo")).lower()
        c_t   = tipo_c.get(tipo, GRIS_OSCURO)
        c_tbg = tipo_bg.get(tipo, GRIS_CLARO)
        prio  = acc.get("prioridad", i + 1)

        _rect(slide, Inches(0.2), y, ANCHO - Inches(0.4), Inches(0.7), bg)

        _rect(slide, Inches(0.22), y + Pt(7), Inches(0.44), Inches(0.38), c_t)
        _texto(slide, str(prio), Inches(0.22), y + Pt(6), Inches(0.44), Pt(28),
               size=12, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

        accion  = str(acc.get("accion", ""))[:105]
        impacto = str(acc.get("impacto_estimado") or "").strip()[:40] or "–"
        plazo   = f"{acc.get('plazo_dias', '?')}d"

        _texto(slide, accion, Inches(0.82), y + Pt(8), Inches(7.0), Inches(0.55),
               size=10, bold=True, color=NEGRO)
        _rect(slide, Inches(8.0), y + Pt(10), Inches(1.6), Pt(22), c_tbg)
        _texto(slide, tipo.upper(), Inches(8.0), y + Pt(9), Inches(1.6), Pt(26),
               size=8, bold=True, color=c_t, align=PP_ALIGN.CENTER)
        _texto(slide, plazo, Inches(9.75), y + Pt(9), Inches(1.0), Inches(0.4),
               size=11, bold=True, color=NEGRO, align=PP_ALIGN.CENTER)
        _texto(slide, impacto, Inches(10.9), y + Pt(9), Inches(2.1), Inches(0.45),
               size=8, color=GRIS)

    _footer(slide, nombre)


# ── Slide 12: Timeline de ejecución ──────────────────────────────────────────

def _slide_timeline(prs, plan: list, nombre: str):
    slide = _slide_en_blanco(prs)
    _rect(slide, 0, 0, ANCHO, ALTO, BLANCO)
    _header(slide, "Timeline de Ejecución",
            "Secuencia temporal recomendada de las acciones priorizadas")

    acciones = sorted(plan, key=lambda x: x.get("prioridad", 99))[:6]
    tipo_c = {"nuevo": VERDE, "renovacion": NARANJA, "upselling": AZUL}

    # Línea de tiempo horizontal
    linea_y = Inches(3.8)
    linea_x_ini = Inches(0.8)
    linea_x_fin = ANCHO - Inches(0.8)
    _linea(slide, linea_x_ini, linea_y, linea_x_fin - linea_x_ini, Pt(3), GRIS_MEDIO)

    total = len(acciones)
    paso  = (linea_x_fin - linea_x_ini) / max(total, 1)

    for i, acc in enumerate(acciones):
        cx     = linea_x_ini + paso * i + paso / 2
        tipo   = str(acc.get("tipo", "nuevo")).lower()
        c      = tipo_c.get(tipo, GRIS_OSCURO)
        plazo  = f"Día {acc.get('plazo_dias', '?')}"
        accion = str(acc.get("accion", ""))[:55]

        # Círculo en línea
        _rect(slide, cx - Inches(0.28), linea_y - Inches(0.28),
              Inches(0.56), Inches(0.56), c)
        _texto(slide, str(i + 1), cx - Inches(0.28), linea_y - Inches(0.3),
               Inches(0.56), Inches(0.56), size=14, bold=True, color=BLANCO,
               align=PP_ALIGN.CENTER)

        # Texto encima / debajo alternados
        if i % 2 == 0:
            _rect(slide, cx - Inches(1.0), Inches(1.2), Inches(2.0), Inches(2.35), GRIS_CLARO)
            _linea(slide, cx - Inches(1.0), Inches(1.2), Inches(2.0), Pt(3), c)
            _texto(slide, plazo, cx - Inches(0.9), Inches(1.3), Inches(1.8), Inches(0.35),
                   size=9, bold=True, color=c, align=PP_ALIGN.CENTER)
            _texto(slide, accion, cx - Inches(0.9), Inches(1.62), Inches(1.8), Inches(1.8),
                   size=9, color=NEGRO, align=PP_ALIGN.CENTER)
            # Conector
            _linea(slide, cx - Pt(1), Inches(3.5), Pt(2), linea_y - Inches(3.5), GRIS_MEDIO)
        else:
            _rect(slide, cx - Inches(1.0), Inches(4.25), Inches(2.0), Inches(2.35), GRIS_CLARO)
            _linea(slide, cx - Inches(1.0), Inches(4.25), Inches(2.0), Pt(3), c)
            _texto(slide, plazo, cx - Inches(0.9), Inches(4.35), Inches(1.8), Inches(0.35),
                   size=9, bold=True, color=c, align=PP_ALIGN.CENTER)
            _texto(slide, accion, cx - Inches(0.9), Inches(4.67), Inches(1.8), Inches(1.8),
                   size=9, color=NEGRO, align=PP_ALIGN.CENTER)
            _linea(slide, cx - Pt(1), linea_y + Pt(3), Pt(2), Inches(4.25) - linea_y - Pt(3), GRIS_MEDIO)

    _footer(slide, nombre)


# ── Slide 13: Argumentario comercial ─────────────────────────────────────────

def _slide_argumentario(prs, argumentario_general: str, productos: list, nombre: str):
    slide = _slide_en_blanco(prs)
    _rect(slide, 0, 0, ANCHO, ALTO, BLANCO)
    _header(slide, "Argumentario Comercial",
            "Mensajes clave adaptados para la conversación con el cliente")

    # Columna izquierda: argumentario general (primeros párrafos)
    _rect(slide, Inches(0.2), Inches(1.2), Inches(6.1), ALTO - Inches(1.4), AZUL_BG)
    _texto(slide, "MENSAJES CLAVE GENERALES", Inches(0.35), Inches(1.3), Inches(5.8), Inches(0.38),
           size=9, bold=True, color=AZUL_OSCURO)
    _linea(slide, Inches(0.35), Inches(1.65), Inches(5.8), Pt(1.5), AZUL)

    if argumentario_general:
        lineas = [l.strip() for l in argumentario_general.split("\n") if l.strip()]
        puntos = [l for l in lineas if not l.startswith("#")][:8]
        for i, punto in enumerate(puntos):
            y = Inches(1.8 + i * 0.65)
            _texto(slide, "▸", Inches(0.35), y, Inches(0.28), Inches(0.5),
                   size=11, bold=True, color=AZUL)
            txt = punto.lstrip("-•▸").strip()[:100]
            _texto(slide, txt, Inches(0.65), y, Inches(5.5), Inches(0.58),
                   size=9, color=NEGRO)
    else:
        _texto(slide, "Argumentario general no disponible.",
               Inches(0.4), Inches(2.2), Inches(5.7), Inches(0.5), size=10, color=GRIS, italic=True)

    # Columna derecha: argumentario por producto top 4
    _rect(slide, Inches(6.5), Inches(1.2), Inches(6.6), ALTO - Inches(1.4), GRIS_CLARO)
    _texto(slide, "ARGUMENTARIO POR PRODUCTO", Inches(6.65), Inches(1.3), Inches(6.3), Inches(0.38),
           size=9, bold=True, color=ROJO_SGS)
    _linea(slide, Inches(6.65), Inches(1.65), Inches(6.3), Pt(1.5), ROJO_SGS)

    prods_sorted = sorted(productos, key=lambda x: -x.get("score_fit", 0))[:4]
    for i, prod in enumerate(prods_sorted):
        y     = Inches(1.8 + i * 1.32)
        score = float(prod.get("score_fit", 0))
        c_s   = _color_score(score)
        _rect(slide, Inches(6.55), y, Inches(6.5), Inches(1.22), BLANCO if i % 2 == 0 else GRIS_CLARO)
        _rect(slide, Inches(6.55), y, Inches(0.12), Inches(1.22), c_s)
        _texto(slide, f"{prod.get('producto', '')[:35]} — {score:.0f}%",
               Inches(6.75), y + Pt(5), Inches(6.1), Inches(0.38),
               size=10, bold=True, color=NEGRO)
        arg = str(prod.get("argumentario", ""))[:130]
        _texto(slide, arg, Inches(6.75), y + Pt(28), Inches(6.1), Inches(0.72),
               size=8.5, color=GRIS_OSCURO)

    _footer(slide, nombre)


# ── Slide 14: Próximos pasos ──────────────────────────────────────────────────

def _slide_proximos_pasos(prs, nombre: str, primera_accion: str):
    slide = _slide_en_blanco(prs)
    _rect(slide, 0, 0, ANCHO, ALTO, NEGRO_SUAVE)
    _rect(slide, 0, 0, Inches(0.08), ALTO, ROJO_SGS)
    _rect(slide, 0, 0, ANCHO, Inches(0.07), ROJO_SGS)

    _texto(slide, "Próximos Pasos", Inches(0.5), Inches(0.4), Inches(10), Inches(0.9),
           size=34, bold=True, color=BLANCO)
    _texto(slide, f"Hoja de ruta para avanzar con {nombre}",
           Inches(0.5), Inches(1.25), Inches(10), Inches(0.45),
           size=14, color=GRIS, italic=True)
    _linea(slide, Inches(0.5), Inches(1.8), Inches(5.5), Pt(2), ROJO_SGS)

    pasos = [
        (primera_accion[:80], "Esta semana", ROJO_SGS),
        ("Presentar propuesta formal y resolver preguntas técnicas del cliente", "Semana 1–2", NARANJA_CLARO),
        ("Ajustar oferta según feedback y negociar condiciones finales", "Semana 2–4", NARANJA_CLARO),
        ("Formalizar acuerdo comercial y arrancar onboarding SGS", "Mes 1–2", VERDE_CLARO),
    ]

    for i, (paso, plazo, color) in enumerate(pasos):
        y = Inches(2.1 + i * 1.22)
        _rect(slide, Inches(0.5), y, Inches(0.6), Inches(0.6), color)
        _texto(slide, str(i + 1), Inches(0.5), y - Pt(1), Inches(0.6), Inches(0.62),
               size=16, bold=True,
               color=NEGRO if color not in (ROJO_SGS,) else BLANCO,
               align=PP_ALIGN.CENTER)
        _rect(slide, Inches(1.35), y + Pt(4), Inches(1.6), Pt(18),
              RGBColor(0x28, 0x28, 0x45))
        _texto(slide, plazo, Inches(1.35), y + Pt(3), Inches(1.6), Pt(22),
               size=8, bold=True, color=color, align=PP_ALIGN.CENTER)
        _texto(slide, paso, Inches(3.15), y, Inches(9.8), Inches(0.55),
               size=14, bold=True, color=BLANCO)

    _footer(slide)


# ── Slide 15: Cierre ──────────────────────────────────────────────────────────

def _slide_cierre(prs, nombre: str):
    slide = _slide_en_blanco(prs)
    _rect(slide, 0, 0, ANCHO, ALTO, NEGRO_SUAVE)
    _rect(slide, 0, 0, ANCHO, Inches(0.08), ROJO_SGS)
    _rect(slide, 0, ALTO - Inches(0.08), ANCHO, Inches(0.08), ROJO_SGS)
    _rect(slide, 0, 0, Inches(0.08), ALTO, ROJO_SGS)
    _rect(slide, ANCHO - Inches(0.08), 0, Inches(0.08), ALTO, ROJO_SGS)

    # Decoración central
    _rect(slide, Inches(3.5), Inches(1.8), Inches(6.3), Inches(4.0),
          RGBColor(0x22, 0x22, 0x40))

    _texto(slide, "SGS España", Inches(4.0), Inches(2.1), Inches(5.3), Inches(1.0),
           size=42, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    _texto(slide, "Inteligencia Comercial", Inches(4.0), Inches(3.0), Inches(5.3), Inches(0.5),
           size=15, color=GRIS, align=PP_ALIGN.CENTER, italic=True)
    _linea(slide, Inches(4.8), Inches(3.55), Inches(3.7), Pt(1.5), ROJO_SGS)
    _texto(slide, f"Propuesta confidencial para {nombre}",
           Inches(4.0), Inches(3.7), Inches(5.3), Inches(0.45),
           size=12, color=RGBColor(0xAA, 0xAA, 0xBB), align=PP_ALIGN.CENTER)
    _texto(slide, date.today().strftime("%B de %Y").capitalize(),
           Inches(4.0), Inches(4.2), Inches(5.3), Inches(0.4),
           size=11, color=GRIS, align=PP_ALIGN.CENTER)

    _texto(slide, "Documento estrictamente confidencial · Uso exclusivo del receptor · No distribuir",
           Inches(1.5), Inches(7.0), Inches(10.3), Inches(0.35),
           size=8, color=GRIS, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# FUNCIÓN PRINCIPAL
# ══════════════════════════════════════════════════════════════════════════════

async def generar_pptx_cuenta(
    cuenta_id: UUID,
    conexion: asyncpg.Connection,
) -> tuple[bytes, str]:
    """Genera el deck PPTX completo (15 slides) para una cuenta. Devuelve (bytes, nombre_fichero)."""

    fila = await conexion.fetchrow(
        """
        SELECT
            c.nombre AS nombre_cuenta,
            pc.productos_recomendados, pc.escenario_optimista,
            pc.escenario_medio, pc.escenario_pesimista, pc.plan_de_accion,
            pc.argumentario_general,
            ie.sector, ie.num_empleados, ie.facturacion_estimada,
            ie.certificaciones_actuales, ie.pain_points, ie.presencia_web,
            ie.oportunidades_detectadas, ie.modelo_usado,
            ie.completado_en AS investigacion_fecha
        FROM propuestas_comerciales pc
        JOIN cuentas c ON c.id = pc.cuenta_id
        LEFT JOIN investigaciones_empresa ie ON ie.id = pc.investigacion_id
        WHERE pc.cuenta_id = $1 AND pc.estado = 'completada'
        ORDER BY pc.creado_en DESC
        LIMIT 1
        """,
        cuenta_id,
    )
    if not fila:
        raise ValueError(f"Sin propuesta completada para la cuenta {cuenta_id}")

    # Pipeline por etapa
    ops_rows = await conexion.fetch(
        """
        SELECT
            COALESCE(etapa::TEXT, 'sin_etapa') AS etapa,
            COUNT(*) AS total_ops,
            COALESCE(SUM(importe), 0) AS importe_total
        FROM oportunidades
        WHERE cuenta_id = $1 AND eliminado_en IS NULL
          AND etapa::TEXT NOT IN ('closed_won','closed_lost','closed_withdrawn')
        GROUP BY etapa
        ORDER BY importe_total DESC
        """,
        cuenta_id,
    )

    # Pipeline por producto
    prods_pipeline = await conexion.fetch(
        """
        SELECT
            COALESCE(p.nombre, o.linea_negocio, 'Sin categoría') AS producto,
            COUNT(o.id) AS total,
            COALESCE(SUM(o.importe), 0) AS importe
        FROM oportunidades o
        LEFT JOIN productos p ON p.id = o.producto_id
        WHERE o.cuenta_id = $1 AND o.eliminado_en IS NULL
          AND o.etapa NOT IN ('closed_won','closed_lost','closed_withdrawn')
        GROUP BY COALESCE(p.nombre, o.linea_negocio, 'Sin categoría')
        ORDER BY importe DESC
        LIMIT 6
        """,
        cuenta_id,
    )

    # Parse datos
    nombre    = fila["nombre_cuenta"]
    opt       = _parse_jsonb(fila["escenario_optimista"], {})
    med       = _parse_jsonb(fila["escenario_medio"], {})
    pes       = _parse_jsonb(fila["escenario_pesimista"], {})
    productos = _parse_jsonb(fila["productos_recomendados"], [])
    plan      = _parse_jsonb(fila["plan_de_accion"], [])
    pain      = _parse_jsonb(fila["pain_points"], [])
    op_ia     = _parse_jsonb(fila["oportunidades_detectadas"], [])
    arg_gral  = fila["argumentario_general"] or ""

    total_pipeline = sum(float(r["importe_total"]) for r in ops_rows)
    total_ops      = sum(r["total_ops"] for r in ops_rows)

    prs = _nueva_presentacion()

    _slide_portada(prs, nombre, fila["sector"], opt, med, pes, len(productos), total_ops, total_pipeline)
    _slide_agenda(prs, nombre)
    _slide_ficha(prs, dict(fila), nombre)
    _slide_diagnostico(prs, pain, op_ia, nombre)
    _slide_pipeline(prs, list(ops_rows), total_ops, total_pipeline, nombre)
    _slide_distribucion_productos(prs, list(prods_pipeline), total_pipeline, nombre)
    _slide_fit_productos(prs, productos, nombre)
    _slide_fit_chart(prs, productos, nombre)
    _slide_escenarios(prs, opt, med, pes, nombre)
    _slide_roi(prs, opt, med, pes, len(productos), nombre)
    _slide_plan_accion(prs, plan, nombre)
    _slide_timeline(prs, plan, nombre)
    _slide_argumentario(prs, arg_gral, productos, nombre)
    _slide_proximos_pasos(prs, nombre, plan[0].get("accion", "Contactar con el cliente") if plan else "Contactar con el cliente")
    _slide_cierre(prs, nombre)

    nombre_archivo = f"deck_{str(cuenta_id)[:8]}_{date.today().isoformat()}.pptx"
    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue(), nombre_archivo
