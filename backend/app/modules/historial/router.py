"""
Endpoints del módulo Historial de documentos generados por agentes IA.
"""

from uuid import UUID

import asyncpg
from fastapi import APIRouter, Depends, HTTPException, Query, status
from fastapi.responses import Response

from app.auth.dependencies import UsuarioAutenticado, obtener_usuario_actual
from app.database import obtener_conexion
from app.modules.historial import servicio
from app.modules.ia.contexto import normalizar_tipo_contexto

router = APIRouter(prefix="/historial", tags=["historial"])


def _content_type_historial(tipo: str) -> str:
    return {
        "pdf": "application/pdf",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "briefing": "text/plain; charset=utf-8",
        "audio": "audio/mpeg",
        "investigacion": "application/json",
        "propuesta": "application/json",
    }.get(tipo, "application/octet-stream")


@router.get("")
async def listar(
    cuenta_id: UUID | None = Query(default=None),
    contexto_tipo: str | None = Query(default=None),
    contexto_id: UUID | None = Query(default=None),
    pagina: int = Query(default=1, ge=1),
    por_pagina: int = Query(default=100, ge=1, le=200),
    sort_by: str = Query(default="creado_en"),
    sort_dir: str = Query(default="desc"),
    usuario: UsuarioAutenticado = Depends(obtener_usuario_actual),
    conexion: asyncpg.Connection = Depends(obtener_conexion),
):
    tipo_normalizado = None
    if contexto_tipo:
        try:
            tipo_normalizado = normalizar_tipo_contexto(contexto_tipo)
        except ValueError as exc:
            raise HTTPException(status_code=422, detail=str(exc))

    return await servicio.listar_historial(
        conexion,
        usuario_id=UUID(usuario.usuario_id),
        es_admin=not usuario.es_comercial,
        cuenta_id=cuenta_id,
        contexto_tipo=tipo_normalizado,
        contexto_id=contexto_id,
        pagina=pagina,
        por_pagina=por_pagina,
        sort_by=sort_by,
        sort_dir=sort_dir,
    )


@router.get("/artefactos")
async def listar_artefactos(
    contexto_tipo: str | None = Query(default=None),
    contexto_id: UUID | None = Query(default=None),
    pagina: int = Query(default=1, ge=1),
    por_pagina: int = Query(default=50, ge=1, le=200),
    sort_by: str = Query(default="creado_en"),
    sort_dir: str = Query(default="desc"),
    usuario: UsuarioAutenticado = Depends(obtener_usuario_actual),
    conexion: asyncpg.Connection = Depends(obtener_conexion),
):
    tipo_normalizado = None
    if contexto_tipo:
        try:
            tipo_normalizado = normalizar_tipo_contexto(contexto_tipo)
        except ValueError as exc:
            raise HTTPException(status_code=422, detail=str(exc))

    return await servicio.listar_artefactos_ia(
        conexion=conexion,
        usuario_id=UUID(usuario.usuario_id),
        es_admin=not usuario.es_comercial,
        contexto_tipo=tipo_normalizado,
        contexto_id=contexto_id,
        pagina=pagina,
        por_pagina=por_pagina,
        sort_by=sort_by,
        sort_dir=sort_dir,
    )


@router.get("/{doc_id}/url")
async def url_descarga(
    doc_id: UUID,
    usuario: UsuarioAutenticado = Depends(obtener_usuario_actual),
    conexion: asyncpg.Connection = Depends(obtener_conexion),
):
    """URL prefirmada de descarga (válida 1 hora)."""
    url = await servicio.obtener_url_descarga(
        conexion,
        doc_id=doc_id,
        usuario_id=UUID(usuario.usuario_id),
        es_admin=not usuario.es_comercial,
    )
    if not url:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Documento no encontrado o sin permisos.")
    return {"url": url}


@router.get("/{doc_id}/descargar")
async def descargar(
    doc_id: UUID,
    usuario: UsuarioAutenticado = Depends(obtener_usuario_actual),
    conexion: asyncpg.Connection = Depends(obtener_conexion),
):
    doc = await servicio.obtener_documento_descarga(
        conexion,
        doc_id=doc_id,
        usuario_id=UUID(usuario.usuario_id),
        es_admin=not usuario.es_comercial,
    )
    if not doc:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Documento no encontrado o sin permisos.")
    headers = {
        "Content-Disposition": f'inline; filename="{doc["nombre_fichero"]}"',
    }
    return Response(
        content=doc["contenido"],
        media_type=_content_type_historial(str(doc["tipo"])),
        headers=headers,
    )


@router.get("/{doc_id}/slides")
async def slides_pptx(
    doc_id: UUID,
    usuario: UsuarioAutenticado = Depends(obtener_usuario_actual),
    conexion: asyncpg.Connection = Depends(obtener_conexion),
):
    """
    Extrae el contenido de cada diapositiva de un PPTX almacenado en MinIO.
    Devuelve: [{index, titulo, cuerpo, notas}]

    Las diapositivas SGS se generan con add_textbox (sin placeholders estándar),
    por eso ordenamos los shapes por posición vertical (shape.top) para identificar
    el título (primer textbox de arriba a abajo).
    """
    fila = await conexion.fetchrow(
        """
        SELECT storage_key, tipo FROM historial_documentos
        WHERE id = $1 AND (usuario_id = $2 OR $3 IS TRUE)
        """,
        doc_id, UUID(usuario.usuario_id), not usuario.es_comercial,
    )
    if not fila:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Documento no encontrado.")
    if fila["tipo"] != "pptx":
        raise HTTPException(status_code=status.HTTP_422_UNPROCESSABLE_ENTITY, detail="El documento no es una presentación PPTX.")

    from app import storage
    contenido = await storage.descargar_fichero(fila["storage_key"])

    import io
    from pptx import Presentation

    prs = Presentation(io.BytesIO(contenido))
    slides = []
    for i, slide in enumerate(prs.slides):
        # Recoger todos los textboxes con texto y ordenar por posición vertical → horizontal
        formas_texto: list[tuple[int, int, str]] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            texto = shape.text_frame.text.strip()
            if not texto:
                continue
            formas_texto.append((shape.top, shape.left, texto))

        formas_texto.sort(key=lambda x: (x[0], x[1]))

        titulo = formas_texto[0][2] if formas_texto else ""
        cuerpo_partes = [t for _, _, t in formas_texto[1:]]

        notas = ""
        try:
            if slide.has_notes_slide:
                ntf = slide.notes_slide.notes_text_frame
                if ntf:
                    notas = ntf.text.strip()
        except Exception:
            pass

        slides.append({
            "index": i + 1,
            "titulo": titulo,
            "cuerpo": "\n".join(cuerpo_partes),
            "notas": notas,
        })

    return {"total": len(slides), "slides": slides}


@router.get("/{doc_id}/texto")
async def texto_documento(
    doc_id: UUID,
    usuario: UsuarioAutenticado = Depends(obtener_usuario_actual),
    conexion: asyncpg.Connection = Depends(obtener_conexion),
):
    """
    Devuelve el contenido de texto de un documento (briefing, investigación, propuesta).
    """
    fila = await conexion.fetchrow(
        """
        SELECT storage_key, tipo FROM historial_documentos
        WHERE id = $1 AND (usuario_id = $2 OR $3 = true)
        """,
        doc_id, UUID(usuario.usuario_id), not usuario.es_comercial,
    )
    if not fila:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Documento no encontrado.")

    from app import storage
    contenido = await storage.descargar_fichero(fila["storage_key"])

    if fila["tipo"] in ("briefing",):
        return {"tipo": fila["tipo"], "texto": contenido.decode("utf-8", errors="replace")}

    # JSON types
    import json as _json
    try:
        datos = _json.loads(contenido)
        return {"tipo": fila["tipo"], "datos": datos}
    except Exception:
        return {"tipo": fila["tipo"], "texto": contenido.decode("utf-8", errors="replace")}


@router.post("/{doc_id}/compartir")
async def compartir(
    doc_id: UUID,
    dias: int = Query(default=7, ge=1, le=30),
    usuario: UsuarioAutenticado = Depends(obtener_usuario_actual),
    conexion: asyncpg.Connection = Depends(obtener_conexion),
):
    """Genera un token público de compartición válido N días."""
    resultado = await servicio.crear_token_comparticion(
        conexion,
        doc_id=doc_id,
        usuario_id=UUID(usuario.usuario_id),
        es_admin=not usuario.es_comercial,
        dias_expiracion=dias,
    )
    if not resultado:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Documento no encontrado o sin permisos.")
    return resultado


@router.delete("/{doc_id}", status_code=status.HTTP_204_NO_CONTENT)
async def eliminar(
    doc_id: UUID,
    usuario: UsuarioAutenticado = Depends(obtener_usuario_actual),
    conexion: asyncpg.Connection = Depends(obtener_conexion),
):
    eliminado = await servicio.eliminar_documento(
        conexion,
        doc_id=doc_id,
        usuario_id=UUID(usuario.usuario_id),
        es_admin=not usuario.es_comercial,
    )
    if not eliminado:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Documento no encontrado o sin permisos.")


# =============================================================================
# Router público — sin autenticación (tokens de compartición)
# =============================================================================

router_compartir_publico = APIRouter(prefix="/s", tags=["compartir"])


@router_compartir_publico.get("/{token}")
async def ver_compartido(
    token: str,
    conexion: asyncpg.Connection = Depends(obtener_conexion),
):
    """
    Endpoint público. Valida el token y sirve el fichero desde MinIO.
    No requiere autenticación — el token actúa como credencial temporal.
    """
    doc = await servicio.obtener_documento_compartido(conexion, token)
    if not doc:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Enlace no encontrado o expirado.")

    content_type = {
        "pdf": "application/pdf",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "audio": "audio/mpeg",
        "briefing": "text/plain; charset=utf-8",
        "investigacion": "application/json",
        "propuesta": "application/json",
    }.get(doc["tipo"], "application/octet-stream")

    return Response(
        content=doc["contenido"],
        media_type=content_type,
        headers={
            "Content-Disposition": f'inline; filename="{doc["nombre_fichero"]}"',
            "Cache-Control": "private, max-age=300",
        },
    )
