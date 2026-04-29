"""
Lógica de generación de decks de visita comercial.
1. IA (Ollama) genera el contenido estructurado de cada slide.
2. python-pptx construye el archivo .pptx con diseño SGS.
3. Al completar, registra artefacto en ia_artefactos vinculado a la cuenta.
"""

import json
import uuid
from pathlib import Path
from uuid import UUID

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

from app.config import configuracion
from app.modules.decks.schemas import SolicitudDeck, TipoDeck, EstadoJob

# =============================================================================
# Constantes de diseño SGS
# =============================================================================

ROJO_SGS = RGBColor(0xC0, 0x00, 0x1A)
NEGRO = RGBColor(0x1A, 0x1A, 0x1A)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
GRIS_CLARO = RGBColor(0xF5, 0xF5, 0xF5)
GRIS_TEXTO = RGBColor(0x66, 0x66, 0x66)

ANCHO = Inches(13.33)   # 16:9 widescreen
ALTO = Inches(7.5)

TMP_DIR = Path(__file__).parent.parent.parent.parent / "tmp_decks"
TMP_DIR.mkdir(exist_ok=True)

# =============================================================================
# Estado de jobs en memoria
# =============================================================================

_jobs: dict[str, dict] = {}


def obtener_estado_job(job_id: str) -> dict | None:
    return _jobs.get(job_id)


def _actualizar_job(job_id: str, **kwargs) -> None:
    if job_id in _jobs:
        _jobs[job_id].update(kwargs)


# =============================================================================
# Prompts por tipo de deck
# =============================================================================

ESTRUCTURA_POR_TIPO = {
    TipoDeck.primera_visita: [
        "Portada",
        "Agenda",
        "SGS en cifras",
        "Propuesta de valor para {sector}",
        "La norma {norma}: qué es y por qué ahora",
        "Beneficios específicos para {empresa}",
        "Proceso de certificación: pasos y plazos",
        "Casos de éxito en {sector}",
        "Inversión y ROI estimado",
        "Próximos pasos",
    ],
    TipoDeck.seguimiento_oferta: [
        "Portada",
        "Resumen de nuestra propuesta",
        "Estado actual del proceso",
        "Respuesta a sus preguntas",
        "Condiciones y garantías",
        "Timeline de implementación",
        "Próximos pasos para cerrar",
    ],
    TipoDeck.upselling: [
        "Portada",
        "Su trayectoria con SGS",
        "Nueva oportunidad detectada",
        "La norma {norma}: complemento natural",
        "Beneficios de la certificación adicional",
        "Condiciones especiales para clientes SGS",
        "Próximos pasos",
    ],
    TipoDeck.propuesta_tecnica: [
        "Portada",
        "Alcance de la propuesta",
        "Descripción de la norma {norma}",
        "Proceso de auditoría",
        "Plazos y calendario",
        "Nuestro equipo auditor",
        "Inversión total",
        "Próximos pasos",
    ],
}

NOMBRES_TIPO = {
    TipoDeck.primera_visita: "Primera Visita",
    TipoDeck.seguimiento_oferta: "Seguimiento de Oferta",
    TipoDeck.upselling: "Upselling / Cross-selling",
    TipoDeck.propuesta_tecnica: "Propuesta Técnica",
}


# =============================================================================
# Generación de contenido con Claude Sonnet
# =============================================================================

async def _generar_contenido_ia(
    solicitud: SolicitudDeck,
    titulos_slides: list[str],
    catalogo_context: str = "",
    matriz_context: str = "",
    ia_config: dict | None = None,
) -> list[dict]:
    """
    Genera el contenido JSON para cada slide usando Ollama.
    Devuelve lista de dicts: [{titulo, puntos: [str], nota_presentador: str}]
    """
    from app.modules.ia.servicio import ConfigIA, llamar_ia
    config = ConfigIA(**ia_config) if ia_config else ConfigIA()

    slides_str = "\n".join(f"{i+1}. {t}" for i, t in enumerate(titulos_slides))

    catalogo_block = f"\nINFORMACIÓN DEL CATÁLOGO SGS (usa estos datos reales):\n{catalogo_context}" if catalogo_context else ""
    matriz_block = f"\nPAIN POINTS Y CERTIFICACIONES DEL SECTOR '{solicitud.sector}':\n{matriz_context}" if matriz_context else ""

    prompt = f"""Eres un consultor comercial senior de SGS España, empresa líder en certificación y testing.
Genera el contenido completo para un deck de visita comercial de tipo "{NOMBRES_TIPO[solicitud.tipo]}".

DATOS DE LA VISITA:
- Empresa cliente: {solicitud.empresa}
- Sector: {solicitud.sector}
- Norma/producto: {solicitud.norma}
- Objetivo de la visita: {solicitud.objetivo_visita}
- Notas adicionales: {solicitud.notas_adicionales or "Ninguna"}
{catalogo_block}{matriz_block}

SLIDES A GENERAR:
{slides_str}

INSTRUCCIONES:
- Para cada slide devuelve exactamente: título, array de 3-4 puntos de contenido concisos (máx 12 palabras cada uno), y una nota del presentador (1-2 frases para guiar al comercial)
- El tono debe ser profesional, consultivo y orientado al valor para el cliente
- Usa los datos reales del catálogo SGS y los pain points del sector en los slides correspondientes
- Usa datos reales de SGS cuando sea relevante (presencia en 140+ países, 100.000+ empleados, fundada en 1878)
- Para la portada: los puntos son [nombre empresa, fecha sugerida, comercial responsable placeholder]
- Responde ÚNICAMENTE con JSON válido, sin texto adicional

FORMATO JSON ESPERADO:
[
  {{
    "titulo": "título del slide",
    "puntos": ["punto 1", "punto 2", "punto 3"],
    "nota_presentador": "texto guía para el comercial"
  }}
]"""

    texto = (await llamar_ia(
        mensajes=[{"role": "user", "content": prompt}],
        system="",
        config=config,
        max_tokens=1200,
    )).strip()
    # Limpiar markdown si Claude lo añade
    if texto.startswith("```"):
        texto = texto.split("```")[1]
        if texto.startswith("json"):
            texto = texto[4:]
    texto = texto.strip()

    return json.loads(texto)


# =============================================================================
# Construcción del PPTX con python-pptx
# =============================================================================

def _añadir_forma_rellena(slide, left, top, width, height, color: RGBColor):
    """Añade un rectángulo de color sólido."""
    from pptx.util import Emu
    forma = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    forma.fill.solid()
    forma.fill.fore_color.rgb = color
    forma.line.fill.background()
    return forma


def _añadir_texto(slide, texto: str, left, top, width, height,
                  tamaño: int, negrita: bool, color: RGBColor,
                  alineacion=PP_ALIGN.LEFT, italica: bool = False):
    """Añade un text box con formato."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alineacion
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(tamaño)
    run.font.bold = negrita
    run.font.italic = italica
    run.font.color.rgb = color
    return txb


def _slide_portada(prs: Presentation, datos: dict, solicitud: SolicitudDeck):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Fondo rojo completo
    _añadir_forma_rellena(slide, 0, 0, ANCHO, ALTO, ROJO_SGS)

    # Banda blanca inferior (1/3 de la altura)
    banda_top = Inches(5.0)
    _añadir_forma_rellena(slide, 0, banda_top, ANCHO, ALTO - banda_top, BLANCO)

    # Logo/marca SGS (texto estilizado)
    _añadir_texto(slide, "SGS", Inches(0.8), Inches(0.5), Inches(3), Inches(1.2),
                  tamaño=60, negrita=True, color=BLANCO)
    _añadir_texto(slide, "España · Inteligencia Comercial",
                  Inches(0.8), Inches(1.6), Inches(6), Inches(0.5),
                  tamaño=14, negrita=False, color=RGBColor(0xFF, 0xCC, 0xCC))

    # Título principal
    titulo = datos.get("titulo", f"Visita a {solicitud.empresa}")
    _añadir_texto(slide, titulo,
                  Inches(0.8), Inches(2.4), Inches(11.5), Inches(1.8),
                  tamaño=36, negrita=True, color=BLANCO)

    # Empresa y norma en banda blanca
    _añadir_texto(slide, solicitud.empresa,
                  Inches(0.8), Inches(5.2), Inches(8), Inches(0.6),
                  tamaño=22, negrita=True, color=ROJO_SGS)
    _añadir_texto(slide, f"{solicitud.norma}  ·  {solicitud.sector}",
                  Inches(0.8), Inches(5.85), Inches(8), Inches(0.5),
                  tamaño=14, negrita=False, color=GRIS_TEXTO)

    # Tipo de presentación
    _añadir_texto(slide, NOMBRES_TIPO[solicitud.tipo],
                  Inches(0.8), Inches(6.5), Inches(8), Inches(0.5),
                  tamaño=12, negrita=False, color=GRIS_TEXTO, italica=True)


def _slide_contenido(prs: Presentation, datos: dict, numero: int, total: int):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Fondo blanco (ya lo es por defecto, solo añadir banda roja lateral)
    _añadir_forma_rellena(slide, 0, 0, Inches(0.12), ALTO, ROJO_SGS)

    # Banda superior sutil
    _añadir_forma_rellena(slide, 0, 0, ANCHO, Inches(0.08), ROJO_SGS)

    # Número de slide (esquina superior derecha)
    _añadir_texto(slide, f"{numero} / {total}",
                  Inches(11.8), Inches(0.15), Inches(1.3), Inches(0.4),
                  tamaño=9, negrita=False, color=GRIS_TEXTO, alineacion=PP_ALIGN.RIGHT)

    # Título del slide
    titulo = datos.get("titulo", "")
    _añadir_texto(slide, titulo,
                  Inches(0.5), Inches(0.3), Inches(12.5), Inches(0.9),
                  tamaño=28, negrita=True, color=ROJO_SGS)

    # Línea separadora
    _añadir_forma_rellena(slide, Inches(0.5), Inches(1.25), Inches(12.33), Inches(0.03), ROJO_SGS)

    # Puntos de contenido
    puntos = datos.get("puntos", [])
    y_inicio = Inches(1.5)
    separacion = Inches(0.85)

    for i, punto in enumerate(puntos[:5]):  # máximo 5 puntos
        y = y_inicio + (i * separacion)
        # Bullet point (círculo rojo)
        _añadir_forma_rellena(slide, Inches(0.5), y + Inches(0.18), Inches(0.12), Inches(0.12), ROJO_SGS)
        # Texto del punto
        _añadir_texto(slide, punto,
                      Inches(0.75), y, Inches(12.2), Inches(0.75),
                      tamaño=18, negrita=False, color=NEGRO)

    # Nota del presentador (zona inferior, tono gris)
    nota = datos.get("nota_presentador", "")
    if nota:
        _añadir_forma_rellena(slide, Inches(0.5), Inches(6.7), Inches(12.33), Inches(0.03),
                              RGBColor(0xDD, 0xDD, 0xDD))
        _añadir_texto(slide, f"💬  {nota}",
                      Inches(0.5), Inches(6.8), Inches(12.33), Inches(0.55),
                      tamaño=10, negrita=False, color=GRIS_TEXTO, italica=True)


def _slide_cierre(prs: Presentation, solicitud: SolicitudDeck):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Fondo rojo
    _añadir_forma_rellena(slide, 0, 0, ANCHO, ALTO, ROJO_SGS)

    _añadir_texto(slide, "Gracias",
                  Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.8),
                  tamaño=60, negrita=True, color=BLANCO)

    _añadir_texto(slide, "¿Alguna pregunta?",
                  Inches(0.8), Inches(3.2), Inches(11.5), Inches(0.8),
                  tamaño=24, negrita=False, color=RGBColor(0xFF, 0xCC, 0xCC))

    _añadir_texto(slide, "sgs.com  ·  SGS España",
                  Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.5),
                  tamaño=14, negrita=False, color=RGBColor(0xFF, 0xCC, 0xCC))


def _construir_pptx(solicitud: SolicitudDeck, contenido: list[dict]) -> Path:
    """Construye el archivo .pptx y lo guarda en TMP_DIR. Devuelve la ruta."""
    prs = Presentation()
    prs.slide_width = ANCHO
    prs.slide_height = ALTO

    total_slides = len(contenido) + 1  # +1 por el slide de cierre

    for i, datos in enumerate(contenido):
        if i == 0:
            _slide_portada(prs, datos, solicitud)
        else:
            _slide_contenido(prs, datos, i, total_slides - 1)

    _slide_cierre(prs, solicitud)

    nombre = f"deck_{solicitud.empresa.replace(' ', '_')[:30]}_{uuid.uuid4().hex[:8]}.pptx"
    ruta = TMP_DIR / nombre
    prs.save(str(ruta))
    return ruta


# =============================================================================
# Tarea principal (ejecutada como BackgroundTask)
# =============================================================================

async def generar_deck_tarea(
    job_id: str,
    solicitud: SolicitudDeck,
    catalogo_context: str = "",
    matriz_context: str = "",
    ia_config: dict | None = None,
    usuario_id: UUID | None = None,
    db_url: str | None = None,
) -> None:
    """
    Genera el deck en background:
    1. Genera contenido con Ollama (enriquecido con catálogo real)
    2. Construye el PPTX
    3. Registra artefacto en ia_artefactos vinculado a la cuenta
    """
    import structlog
    log = structlog.get_logger()

    try:
        # Paso 1: Generar contenido IA
        _actualizar_job(job_id,
                        estado=EstadoJob.generando_contenido,
                        progreso=10,
                        mensaje="Generando contenido con IA...")

        titulos_base = ESTRUCTURA_POR_TIPO[solicitud.tipo]
        titulos = [
            t.format(
                empresa=solicitud.empresa,
                sector=solicitud.sector,
                norma=solicitud.norma,
            )
            for t in titulos_base
        ]

        if solicitud.num_slides < len(titulos):
            titulos = titulos[:solicitud.num_slides]

        _actualizar_job(job_id, progreso=30)

        contenido = await _generar_contenido_ia(solicitud, titulos, catalogo_context, matriz_context, ia_config)

        _actualizar_job(job_id, progreso=60)

        # Paso 2: Construir PPTX
        _actualizar_job(job_id,
                        estado=EstadoJob.construyendo_slides,
                        progreso=65,
                        mensaje="Construyendo presentación...")

        ruta = _construir_pptx(solicitud, contenido)

        # Paso 3: Registrar artefacto en DB vinculado a la cuenta
        if usuario_id and db_url and solicitud.cuenta_id:
            try:
                import asyncpg
                from app.modules.artefactos.servicio import registrar_version_artefacto
                conn = await asyncpg.connect(db_url.replace("+asyncpg", ""))
                try:
                    await registrar_version_artefacto(
                        conn,
                        tipo="deck",
                        subtipo="pptx",
                        entidad_tipo="cuenta",
                        entidad_id=solicitud.cuenta_id,
                        cuenta_id=solicitud.cuenta_id,
                        usuario_id=usuario_id,
                        titulo=f"Deck · {solicitud.empresa} · {NOMBRES_TIPO[solicitud.tipo]}",
                        storage_key=str(ruta),
                        modelo=(ia_config or {}).get("ollama_modelo") or configuracion.OLLAMA_MODEL_DEFAULT,
                        origen_tabla="decks_jobs",
                        origen_id=job_id,
                        metadatos={
                            "tipo_deck": solicitud.tipo,
                            "empresa": solicitud.empresa,
                            "sector": solicitud.sector,
                            "norma": solicitud.norma,
                            "num_slides": solicitud.num_slides,
                            "archivo": ruta.name,
                        },
                    )
                finally:
                    await conn.close()
            except Exception as e:
                log.warning("deck_artefacto_registro_fallido", job_id=job_id, error=str(e))

        _actualizar_job(job_id,
                        estado=EstadoJob.completado,
                        progreso=100,
                        mensaje="Presentación lista para descargar",
                        archivo=ruta.name)

    except Exception as exc:
        _actualizar_job(job_id,
                        estado=EstadoJob.error,
                        progreso=0,
                        mensaje=f"Error: {str(exc)[:200]}")
