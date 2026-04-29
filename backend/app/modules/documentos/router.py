"""
Router del módulo Documentos.
Endpoints: subir, listar, descargar, eliminar.
"""

from uuid import UUID

import asyncpg
from fastapi import APIRouter, Depends, File, Form, HTTPException, Query, UploadFile, status
from fastapi.responses import FileResponse

from app.auth.dependencies import UsuarioAutenticado, obtener_usuario_actual
from app.database import obtener_conexion
from app.modules.documentos import servicio
from app.modules.documentos.schemas import DocumentoRead

router = APIRouter(prefix="/documentos", tags=["documentos"])


@router.post("/subir", response_model=DocumentoRead, status_code=status.HTTP_201_CREATED)
async def subir_documento(
    archivo: UploadFile = File(...),
    oportunidad_id: UUID | None = Form(default=None),
    cuenta_id: UUID | None = Form(default=None),
    usuario: UsuarioAutenticado = Depends(obtener_usuario_actual),
    conexion: asyncpg.Connection = Depends(obtener_conexion),
):
    contenido = await archivo.read()
    try:
        return await servicio.subir_documento(
            conexion=conexion,
            usuario_id=UUID(usuario.usuario_id),
            nombre_original=archivo.filename or "sin_nombre",
            contenido=contenido,
            tipo_mime=archivo.content_type,
            oportunidad_id=oportunidad_id,
            cuenta_id=cuenta_id,
        )
    except ValueError as e:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=str(e))


@router.get("", response_model=dict)
async def listar_mis_documentos(
    oportunidad_id: UUID | None = Query(default=None),
    propietario_id: UUID | None = Query(default=None),
    busqueda: str | None = Query(default=None, max_length=200),
    pagina: int = Query(default=1, ge=1),
    por_pagina: int = Query(default=20, ge=1, le=100),
    sort_by: str = Query(default="creado_en"),
    sort_dir: str = Query(default="desc"),
    usuario: UsuarioAutenticado = Depends(obtener_usuario_actual),
    conexion: asyncpg.Connection = Depends(obtener_conexion),
):
    usuario_objetivo_id = UUID(usuario.usuario_id)
    if propietario_id is not None and usuario.es_manager:
        usuario_objetivo_id = propietario_id

    return await servicio.listar_documentos(
        conexion=conexion,
        usuario_id=usuario_objetivo_id,
        oportunidad_id=oportunidad_id,
        busqueda=busqueda,
        pagina=pagina,
        por_pagina=por_pagina,
        sort_by=sort_by,
        sort_dir=sort_dir,
    )


@router.get("/{documento_id}/descargar")
async def descargar_documento(
    documento_id: UUID,
    usuario: UsuarioAutenticado = Depends(obtener_usuario_actual),
    conexion: asyncpg.Connection = Depends(obtener_conexion),
):
    resultado = await servicio.obtener_ruta_descarga(
        conexion=conexion,
        documento_id=documento_id,
        usuario_id=UUID(usuario.usuario_id),
    )
    if not resultado:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Documento no encontrado.",
        )
    ruta, nombre_original = resultado
    return FileResponse(
        path=str(ruta),
        filename=nombre_original,
        media_type="application/octet-stream",
    )


@router.get("/{documento_id}/slides")
async def slides_pptx_documento(
    documento_id: UUID,
    usuario: UsuarioAutenticado = Depends(obtener_usuario_actual),
    conexion: asyncpg.Connection = Depends(obtener_conexion),
):
    fila = await conexion.fetchrow(
        "SELECT nombre_guardado, nombre_original, tipo_mime FROM documentos WHERE id = $1 AND usuario_id = $2",
        documento_id, UUID(usuario.usuario_id),
    )
    if not fila:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Documento no encontrado.")

    tipo = (fila["tipo_mime"] or "").lower()
    nombre = (fila["nombre_original"] or "").lower()
    if not ("presentation" in tipo or nombre.endswith(".pptx") or nombre.endswith(".ppt")):
        raise HTTPException(status_code=status.HTTP_422_UNPROCESSABLE_ENTITY, detail="El documento no es una presentación PPTX.")

    from app.modules.documentos.servicio import ruta_documento
    ruta = ruta_documento(fila["nombre_guardado"])
    if not ruta.exists():
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Fichero no encontrado en disco.")

    import io
    from pptx import Presentation

    prs = Presentation(io.BytesIO(ruta.read_bytes()))
    slides = []
    for i, slide in enumerate(prs.slides):
        formas_texto: list[tuple[int, int, str]] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            texto = shape.text_frame.text.strip()
            if not texto:
                continue
            formas_texto.append((shape.top, shape.left, texto))
        formas_texto.sort(key=lambda x: (x[0], x[1]))
        titulo = formas_texto[0][2] if formas_texto else ""
        cuerpo_partes = [t for _, _, t in formas_texto[1:]]
        notas = ""
        try:
            if slide.has_notes_slide:
                ntf = slide.notes_slide.notes_text_frame
                if ntf:
                    notas = ntf.text.strip()
        except Exception:
            pass
        slides.append({"index": i + 1, "titulo": titulo, "cuerpo": "\n".join(cuerpo_partes), "notas": notas})

    return {"total": len(slides), "slides": slides}


@router.post("/{documento_id}/compartir")
async def compartir_documento(
    documento_id: UUID,
    dias: int = Query(default=7, ge=1, le=30),
    usuario: UsuarioAutenticado = Depends(obtener_usuario_actual),
    conexion: asyncpg.Connection = Depends(obtener_conexion),
):
    """Genera token público de compartición para un documento de usuario."""
    existe = await conexion.fetchval(
        "SELECT id FROM documentos WHERE id = $1 AND usuario_id = $2",
        documento_id, UUID(usuario.usuario_id),
    )
    if not existe:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Documento no encontrado.")

    fila = await conexion.fetchrow(
        """
        INSERT INTO artefacto_compartidos (documentos_id, creado_por, expira_en)
        VALUES ($1, $2, now() + ($3 || ' days')::INTERVAL)
        RETURNING token, expira_en::TEXT AS expira_en
        """,
        documento_id,
        UUID(usuario.usuario_id),
        str(dias),
    )
    return dict(fila)


@router.post("/{documento_id}/transcribir")
async def transcribir_documento(
    documento_id: UUID,
    usuario: UsuarioAutenticado = Depends(obtener_usuario_actual),
    conexion: asyncpg.Connection = Depends(obtener_conexion),
):
    """Transcribe un archivo de audio subido usando Whisper (local o OpenAI según config)."""
    fila = await conexion.fetchrow(
        "SELECT nombre_guardado, nombre_original, tipo_mime, cuenta_id FROM documentos WHERE id = $1 AND usuario_id = $2",
        documento_id, UUID(usuario.usuario_id),
    )
    if not fila:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Documento no encontrado.")

    ruta = servicio.ruta_documento(fila["nombre_guardado"])
    if not ruta.exists():
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Fichero no encontrado en disco.")

    from app.config import configuracion as cfg

    try:
        if cfg.OPENAI_API_KEY:
            from openai import AsyncOpenAI
            cliente = AsyncOpenAI(api_key=cfg.OPENAI_API_KEY)
            with open(ruta, "rb") as f:
                result = await cliente.audio.transcriptions.create(
                    model="whisper-1",
                    file=f,
                    language="es",
                )
            texto = result.text
        else:
            import asyncio
            import functools

            def _transcribir_sync(ruta_str: str, modelo: str) -> str:
                from faster_whisper import WhisperModel
                model = WhisperModel(modelo, device="cpu", compute_type="int8")
                segments, _ = model.transcribe(ruta_str, language="es", beam_size=1)
                return " ".join(seg.text.strip() for seg in segments)

            loop = asyncio.get_event_loop()
            texto = await loop.run_in_executor(
                None,
                functools.partial(_transcribir_sync, str(ruta), cfg.WHISPER_MODELO),
            )
    except Exception as e:
        raise HTTPException(status_code=status.HTTP_502_BAD_GATEWAY, detail=f"Error en transcripción: {e}")

    # Si el documento está vinculado a una cuenta, guardar también en contenido_extraido
    # para que el pipeline de IA (informes, briefings, presentaciones) pueda usarlo.
    if fila["cuenta_id"]:
        await conexion.execute(
            """
            UPDATE documentos
            SET transcripcion_texto = $1,
                contenido_extraido = COALESCE(contenido_extraido, $1)
            WHERE id = $2
            """,
            texto, documento_id,
        )
    else:
        await conexion.execute(
            "UPDATE documentos SET transcripcion_texto = $1 WHERE id = $2",
            texto, documento_id,
        )

    return {"id": str(documento_id), "transcripcion": texto}


@router.get("/{documento_id}/transcripcion")
async def obtener_transcripcion(
    documento_id: UUID,
    usuario: UsuarioAutenticado = Depends(obtener_usuario_actual),
    conexion: asyncpg.Connection = Depends(obtener_conexion),
):
    """Devuelve la transcripción guardada de un documento de audio."""
    texto = await conexion.fetchval(
        "SELECT transcripcion_texto FROM documentos WHERE id = $1 AND usuario_id = $2",
        documento_id, UUID(usuario.usuario_id),
    )
    if texto is None:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Sin transcripción disponible.")
    return {"id": str(documento_id), "transcripcion": texto}


@router.delete("/{documento_id}", status_code=status.HTTP_204_NO_CONTENT)
async def eliminar_documento(
    documento_id: UUID,
    usuario: UsuarioAutenticado = Depends(obtener_usuario_actual),
    conexion: asyncpg.Connection = Depends(obtener_conexion),
):
    eliminado = await servicio.eliminar_documento(
        conexion=conexion,
        documento_id=documento_id,
        usuario_id=UUID(usuario.usuario_id),
    )
    if not eliminado:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Documento no encontrado.",
        )
